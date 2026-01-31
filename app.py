import streamlit as st
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import requests
from io import BytesIO

# Page config
st.set_page_config(
    page_title="Genis 2.0 - Slideshow Generator",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS
st.markdown("""
<style>
    [data-testid="stSidebar"] {
        min-width: 300px !important;
        max-width: 300px !important;
    }
    
    .block-container {
        max-width: 900px !important;
        padding-left: 2rem;
        padding-right: 2rem;
    }
    
    .big-title {
        text-align: center;
        font-size: 3.5rem;
        font-weight: bold;
        margin-bottom: 0.5rem;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    
    .subtitle {
        text-align: center;
        font-size: 1.3rem;
        color: #666;
        margin-bottom: 2rem;
    }
    
    .stButton > button {
        width: 100%;
        height: 3.5rem;
        font-size: 1.3rem;
        font-weight: bold;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        border-radius: 10px;
        margin-top: 1.5rem;
    }
    
    .stButton > button:hover {
        background: linear-gradient(90deg, #764ba2 0%, #667eea 100%);
        transform: translateY(-2px);
        box-shadow: 0 4px 12px rgba(102, 126, 234, 0.4);
    }
    
    .info-card {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 10px;
        border-left: 4px solid #667eea;
        margin: 1rem 0;
    }
    
    .color-preview {
        display: inline-block;
        width: 20px;
        height: 20px;
        border-radius: 3px;
        margin-right: 8px;
        vertical-align: middle;
        border: 1px solid #ddd;
    }
</style>
""", unsafe_allow_html=True)

# Get API keys from secrets
try:
    groq_key = st.secrets["GROQ_API_KEY"]
    client = Groq(api_key=groq_key)
    api_available = True
except:
    api_available = False

try:
    pexels_key = st.secrets["PEXELS_API_KEY"]
    pexels_available = True
except:
    pexels_available = False

# Initialize session state
if "example_prompts" not in st.session_state:
    st.session_state.example_prompts = []
if "selected_example" not in st.session_state:
    st.session_state.selected_example = ""

# Color themes
COLOR_THEMES = {
    "Blue Ocean 🌊": {
        "dark": RGBColor(25, 50, 100),
        "accent": RGBColor(66, 135, 245),
        "preview": "#1a3264"
    },
    "Forest Green 🌲": {
        "dark": RGBColor(27, 94, 32),
        "accent": RGBColor(102, 187, 106),
        "preview": "#1b5e20"
    },
    "Sunset Orange 🌅": {
        "dark": RGBColor(191, 54, 12),
        "accent": RGBColor(255, 152, 0),
        "preview": "#bf360c"
    },
    "Royal Purple 👑": {
        "dark": RGBColor(74, 20, 140),
        "accent": RGBColor(156, 39, 176),
        "preview": "#4a148c"
    },
    "Ruby Red 💎": {
        "dark": RGBColor(183, 28, 28),
        "accent": RGBColor(244, 67, 54),
        "preview": "#b71c1c"
    },
    "Midnight Black 🌙": {
        "dark": RGBColor(33, 33, 33),
        "accent": RGBColor(97, 97, 97),
        "preview": "#212121"
    },
    "Teal Wave 🏄": {
        "dark": RGBColor(0, 77, 64),
        "accent": RGBColor(0, 150, 136),
        "preview": "#004d40"
    },
    "Pink Blossom 🌸": {
        "dark": RGBColor(136, 14, 79),
        "accent": RGBColor(233, 30, 99),
        "preview": "#880e4f"
    }
}

def get_pexels_image(query):
    """Fetch image from Pexels API"""
    if not pexels_available:
        return None
    
    try:
        headers = {
            "Authorization": st.secrets["PEXELS_API_KEY"]
        }
        url = f"https://api.pexels.com/v1/search?query={query}&per_page=1&orientation=landscape"
        response = requests.get(url, headers=headers, timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            if data.get("photos") and len(data["photos"]) > 0:
                image_url = data["photos"][0]["src"]["large"]
                img_response = requests.get(image_url, timeout=10)
                if img_response.status_code == 200:
                    return BytesIO(img_response.content)
        return None
    except Exception as e:
        return None

# Sidebar
with st.sidebar:
    st.markdown("## 🎯 Genis 2.0")
    st.markdown("---")
    
    st.markdown("### ⚡ Features")
    st.markdown("""
    - 🎨 Custom color themes
    - 🖼️ AI-powered images
    - 📊 Professional layouts
    - 📱 Mobile compatible
    - 💾 Instant download
    """)
    
    st.markdown("---")
    
    st.markdown("### 📱 Device Support")
    st.success("✅ Works on all devices")
    st.markdown("""
    **iPhone/iPad:**  
    PowerPoint or Keynote
    
    **Android:**  
    PowerPoint or Google Slides
    
    **Desktop:**  
    PowerPoint, Google Slides, LibreOffice
    """)
    
    st.markdown("---")
    
    st.markdown("### 🔢 Limits")
    st.info("**Max slides:** 100  \n**Min slides:** 5")
    
    st.markdown("---")
    
    # Status indicators
    if not api_available:
        st.error("⚠️ Groq API not configured")
    else:
        st.success("✅ AI System Ready")
    
    if not pexels_available:
        st.warning("⚠️ Images disabled (no Pexels key)")
    else:
        st.success("✅ Image System Ready")
    
    st.markdown("---")
    st.caption("© 2025 Genis 2.0")

# Main content
st.markdown('<div class="big-title">🎯 Genis 2.0</div>', unsafe_allow_html=True)
st.markdown('<div class="subtitle">Professional Slideshow Generator</div>', unsafe_allow_html=True)

# Generate example prompts
if not st.session_state.example_prompts and api_available:
    with st.spinner("🎲 Generating example ideas..."):
        try:
            example_prompt = """Generate 4 diverse and interesting slideshow topic examples. 
            Make them varied across different fields (business, education, science, technology, health, history, etc.).
            
            Format each as a complete request like:
            "Create a 12-slide presentation about [topic]"
            
            Make them specific and engaging. Return ONLY the 4 examples, one per line, nothing else."""
            
            response = client.chat.completions.create(
                messages=[{"role": "user", "content": example_prompt}],
                model="llama-3.3-70b-versatile",
                temperature=0.9,
                max_tokens=300,
            )
            
            examples = response.choices[0].message.content.strip().split('\n')
            st.session_state.example_prompts = [ex.strip() for ex in examples if ex.strip()]
        except:
            st.session_state.example_prompts = [
                "Create a 10-slide presentation about renewable energy sources",
                "Create a 15-slide presentation about the human brain and memory",
                "Create a 12-slide presentation about digital marketing strategies",
                "Create an 8-slide presentation about ancient Egyptian civilization"
            ]

# Show examples
if st.session_state.example_prompts:
    st.markdown("### 💡 Need inspiration? Try these:")
    
    cols = st.columns(2)
    for idx, example in enumerate(st.session_state.example_prompts[:4]):
        with cols[idx % 2]:
            if st.button(f"📄 {example}", key=f"example_{idx}", use_container_width=True):
                st.session_state.selected_example = example

# Refresh examples button
col1, col2, col3 = st.columns([1, 1, 1])
with col2:
    if st.button("🔄 New Examples", use_container_width=True):
        st.session_state.example_prompts = []
        st.rerun()

st.markdown("---")

# Input fields
st.markdown("### 👤 Your Information")

author_name = st.text_input(
    "Your Name (will appear on title slide)",
    placeholder="e.g., John Smith",
    help="This will be shown as the author on the first slide"
)

st.markdown("### 📝 What would you like to create?")

# Topic input
if st.session_state.selected_example:
    topic = st.text_area(
        "Describe your slideshow:",
        value=st.session_state.selected_example,
        placeholder="Example: Create a 10-slide presentation about climate change",
        height=100,
        label_visibility="collapsed",
        key="topic_input"
    )
else:
    topic = st.text_area(
        "Describe your slideshow:",
        placeholder="Example: Create a 10-slide presentation about climate change",
        height=100,
        label_visibility="collapsed",
        key="topic_input"
    )

# Advanced options
with st.expander("⚙️ Customization Options", expanded=True):
    col1, col2 = st.columns(2)
    
    with col1:
        num_slides = st.number_input(
            "Number of slides", 
            min_value=5, 
            max_value=100, 
            value=10,
            help="Choose between 5 and 100 slides"
        )
        
        style = st.selectbox(
            "Presentation style",
            ["Professional", "Creative", "Educational", "Minimal", "Bold", "Modern"]
        )
    
    with col2:
        # Color theme selector
        st.markdown("**Color Theme**")
        selected_theme = st.selectbox(
            "Choose your color scheme",
            list(COLOR_THEMES.keys()),
            label_visibility="collapsed"
        )
        
        # Show color preview
        theme_colors = COLOR_THEMES[selected_theme]
        st.markdown(f"""
        <div style="margin-top: 10px;">
            <span class="color-preview" style="background-color: {theme_colors['preview']};"></span>
            <span style="font-size: 0.9rem; color: #666;">Selected theme preview</span>
        </div>
        """, unsafe_allow_html=True)
        
        # Image option
        if pexels_available:
            add_images = st.checkbox(
                "🖼️ Add professional images to slides",
                value=True,
                help="Automatically adds relevant images from Pexels"
            )
        else:
            st.info("📸 Images disabled - add PEXELS_API_KEY to enable")
            add_images = False

# Generate button
if st.button("🚀 Generate Slideshow"):
    st.session_state.selected_example = ""
    
    if not api_available:
        st.error("❌ System not configured. Please contact administrator.")
    elif not topic or len(topic.strip()) < 10:
        st.warning("⚠️ Please enter a more detailed topic for your slideshow.")
    elif not author_name or len(author_name.strip()) < 2:
        st.warning("⚠️ Please enter your name to continue.")
    else:
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        try:
            # Step 1: Preparing
            status_text.text("🔄 Initializing Genis 2.0...")
            progress_bar.progress(10)
            
            # Create enhanced prompt
            enhanced_prompt = f"""Create slideshow content about: "{topic}"

Requirements:
- Number of slides: {num_slides}
- Style: {style}

Provide content in this EXACT format:

SLIDE 1:
TITLE: [Compelling title for the presentation]
CONTENT:
- [This is the title slide, no bullet points needed]

SLIDE 2:
TITLE: [Slide title]
CONTENT:
- [Bullet point 1]
- [Bullet point 2]
- [Bullet point 3]

Continue for all {num_slides} slides. Last slide should be a strong conclusion.
Make content informative, engaging, and well-structured."""

            # Step 2: Generating content
            status_text.text("🧠 Generating content...")
            progress_bar.progress(20)
            
            chat_completion = client.chat.completions.create(
                messages=[{"role": "user", "content": enhanced_prompt}],
                model="llama-3.3-70b-versatile",
                temperature=0.7,
                max_tokens=8000,
            )
            
            response_text = chat_completion.choices[0].message.content
            
            # Step 3: Parsing content
            status_text.text("📊 Structuring slides...")
            progress_bar.progress(35)
            
            lines = response_text.split('\n')
            slides_data = []
            current_slide = None
            
            for line in lines:
                line = line.strip()
                if line.startswith('SLIDE '):
                    if current_slide:
                        slides_data.append(current_slide)
                    current_slide = {"title": "", "content": []}
                elif line.startswith('TITLE:'):
                    if current_slide is not None:
                        current_slide["title"] = line.replace('TITLE:', '').strip()
                elif (line.startswith('- ') or line.startswith('• ')) and current_slide is not None:
                    current_slide["content"].append(line[2:].strip())
            
            if current_slide:
                slides_data.append(current_slide)
            
            if not slides_data:
                st.error("❌ Failed to generate slideshow. Please try again.")
                progress_bar.empty()
                status_text.empty()
            else:
                # Step 4: Fetching images
                image_data = {}
                images_fetched = 0
                
                if add_images and pexels_available:
                    status_text.text("🖼️ Finding perfect images...")
                    progress_bar.progress(50)
                    
                    # Get main keywords from topic
                    topic_words = topic.lower().replace("create", "").replace("presentation", "").replace("about", "").replace("slide", "").strip()
                    
                    # Fetch images for select slides (every 3rd slide, max 4 images)
                    for idx in range(2, min(len(slides_data), 14), 3):
                        if idx < len(slides_data):
                            # Use slide title for search
                            search_query = slides_data[idx]["title"].split()[:3]
                            search_term = " ".join(search_query)
                            
                            img = get_pexels_image(search_term)
                            if img:
                                image_data[idx] = img
                                images_fetched += 1
                                status_text.text(f"🖼️ Found {images_fetched} images...")
                
                # Step 5: Creating presentation
                status_text.text("🎨 Designing presentation...")
                progress_bar.progress(70)
                
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625)
                
                # Get selected colors
                DARK_COLOR = COLOR_THEMES[selected_theme]["dark"]
                ACCENT_COLOR = COLOR_THEMES[selected_theme]["accent"]
                WHITE = RGBColor(255, 255, 255)
                GRAY = RGBColor(80, 80, 80)
                
                for idx, slide_data in enumerate(slides_data):
                    blank_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_layout)
                    
                    if idx == 0:
                        # Title slide
                        background = slide.shapes.add_shape(
                            1, 0, 0, prs.slide_width, prs.slide_height
                        )
                        background.fill.solid()
                        background.fill.fore_color.rgb = DARK_COLOR
                        background.line.fill.background()
                        
                        title_box = slide.shapes.add_textbox(
                            Inches(1), Inches(1.5), Inches(8), Inches(1.8)
                        )
                        title_frame = title_box.text_frame
                        title_frame.text = slide_data["title"]
                        title_frame.word_wrap = True
                        title_para = title_frame.paragraphs[0]
                        title_para.font.size = Pt(48)
                        title_para.font.bold = True
                        title_para.font.color.rgb = WHITE
                        title_para.alignment = PP_ALIGN.CENTER
                        
                        subtitle_box = slide.shapes.add_textbox(
                            Inches(1), Inches(3.5), Inches(8), Inches(0.6)
                        )
                        subtitle_frame = subtitle_box.text_frame
                        subtitle_frame.text = f"by {author_name.strip()}"
                        subtitle_para = subtitle_frame.paragraphs[0]
                        subtitle_para.font.size = Pt(24)
                        subtitle_para.font.color.rgb = ACCENT_COLOR
                        subtitle_para.alignment = PP_ALIGN.CENTER
                        
                    else:
                        # Content slides
                        background = slide.shapes.add_shape(
                            1, 0, 0, prs.slide_width, prs.slide_height
                        )
                        background.fill.solid()
                        background.fill.fore_color.rgb = WHITE
                        background.line.fill.background()
                        
                        accent = slide.shapes.add_shape(
                            1, 0, 0, prs.slide_width, Inches(0.2)
                        )
                        accent.fill.solid()
                        accent.fill.fore_color.rgb = ACCENT_COLOR
                        accent.line.fill.background()
                        
                        # Check if slide has image
                        has_image = idx in image_data
                        
                        if has_image:
                            # Split layout
                            title_box = slide.shapes.add_textbox(
                                Inches(0.5), Inches(0.5), Inches(5), Inches(0.8)
                            )
                            title_frame = title_box.text_frame
                            title_frame.text = slide_data["title"]
                            title_para = title_frame.paragraphs[0]
                            title_para.font.size = Pt(30)
                            title_para.font.bold = True
                            title_para.font.color.rgb = DARK_COLOR
                            
                            # Add image
                            try:
                                img_stream = image_data[idx]
                                pic = slide.shapes.add_picture(
                                    img_stream,
                                    Inches(5.5), Inches(1.2),
                                    width=Inches(4), height=Inches(3.8)
                                )
                            except Exception as e:
                                pass
                            
                            # Content
                            if slide_data["content"]:
                                content_box = slide.shapes.add_textbox(
                                    Inches(0.8), Inches(1.7), Inches(4.4), Inches(3.3)
                                )
                                text_frame = content_box.text_frame
                                text_frame.word_wrap = True
                                
                                for point in slide_data["content"]:
                                    p = text_frame.add_paragraph()
                                    p.text = point
                                    p.level = 0
                                    p.font.size = Pt(16)
                                    p.font.color.rgb = GRAY
                                    p.space_before = Pt(10)
                                    p.space_after = Pt(6)
                        else:
                            # Full width
                            title_box = slide.shapes.add_textbox(
                                Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
                            )
                            title_frame = title_box.text_frame
                            title_frame.text = slide_data["title"]
                            title_para = title_frame.paragraphs[0]
                            title_para.font.size = Pt(36)
                            title_para.font.bold = True
                            title_para.font.color.rgb = DARK_COLOR
                            
                            if slide_data["content"]:
                                content_box = slide.shapes.add_textbox(
                                    Inches(1.2), Inches(1.8), Inches(7.6), Inches(3.2)
                                )
                                text_frame = content_box.text_frame
                                text_frame.word_wrap = True
                                
                                for point in slide_data["content"]:
                                    p = text_frame.add_paragraph()
                                    p.text = point
                                    p.level = 0
                                    p.font.size = Pt(20)
                                    p.font.color.rgb = GRAY
                                    p.space_before = Pt(14)
                                    p.space_after = Pt(8)
                
                # Step 6: Finalizing
                status_text.text("✅ Finalizing presentation...")
                progress_bar.progress(95)
                
                pptx_bytes = io.BytesIO()
                prs.save(pptx_bytes)
                pptx_bytes.seek(0)
                
                progress_bar.progress(100)
                progress_bar.empty()
                status_text.empty()
                
                st.balloons()
                st.success(f"🎉 **Success!** Your {len(slides_data)}-slide presentation is ready!")
                
                # Stats
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("Slides", len(slides_data))
                with col2:
                    st.metric("Style", style)
                with col3:
                    st.metric("Theme", selected_theme.split()[0])
                with col4:
                    st.metric("Images", images_fetched)
                
                safe_title = slides_data[0]['title'][:40].replace(' ', '_').replace('/', '_').replace('\\', '_')
                file_name = f"{safe_title}.pptx"
                
                st.download_button(
                    label="📥 Download Presentation",
                    data=pptx_bytes.getvalue(),
                    file_name=file_name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
                
                st.markdown("""
                <div class="info-card">
                    <strong>📱 Opening on your device:</strong><br>
                    • <strong>iPhone/iPad:</strong> PowerPoint or Keynote app<br>
                    • <strong>Android:</strong> PowerPoint or Google Slides<br>
                    • <strong>Desktop:</strong> PowerPoint, Google Slides, or LibreOffice
                </div>
                """, unsafe_allow_html=True)
                
        except Exception as e:
            progress_bar.empty()
            status_text.empty()
            st.error(f"❌ Error: {str(e)}")
            st.info("Please try again or rephrase your topic.")

# Footer
st.markdown("---")
st.markdown("""
<div style="text-align: center; color: #666; padding: 2rem 0;">
    <h3 style="background: linear-gradient(90deg, #667eea 0%, #764ba2 100%); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">
        Genis 2.0
    </h3>
    <p>Professional Slideshows in Seconds</p>
    <p style="font-size: 0.9rem; margin-top: 1rem;">© 2025 Genis 2.0. All presentations created are owned by you.</p>
</div>
""", unsafe_allow_html=True)
